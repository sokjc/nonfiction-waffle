"""PowerPoint (.pptx) loader using python-pptx.

Extracts text from all slides, including shapes, tables, and grouped shapes,
and returns one LangChain ``Document`` per slide.
"""

from __future__ import annotations

import logging
from pathlib import Path

from langchain_core.documents import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class PptxLoader:
    """Load a .pptx file and return one Document per slide.

    Compatible with the LangChain loader interface (``load()`` returns a list
    of ``Document`` objects).
    """

    def __init__(self, file_path: str) -> None:
        self.file_path = file_path

    @staticmethod
    def _extract_shape_text(shape) -> str:
        """Recursively extract text from a shape, including tables and groups."""
        parts: list[str] = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)

        if shape.has_table:
            for row in shape.table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        if shape.shape_type is not None:
            # Group shapes contain child shapes
            try:
                for child in shape.shapes:
                    child_text = PptxLoader._extract_shape_text(child)
                    if child_text:
                        parts.append(child_text)
            except AttributeError:
                pass

        return "\n".join(parts)

    def load(self) -> list[Document]:
        """Load the PowerPoint file and return a list of Documents."""
        prs = Presentation(self.file_path)
        docs: list[Document] = []
        file_name = Path(self.file_path).name

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                text = self._extract_shape_text(shape)
                if text:
                    slide_parts.append(text)

            page_content = "\n\n".join(slide_parts)
            if not page_content.strip():
                continue

            docs.append(
                Document(
                    page_content=page_content,
                    metadata={
                        "source": self.file_path,
                        "slide_number": slide_num,
                        "source_file": file_name,
                    },
                )
            )

        logger.info("Loaded %d slide(s) from %s", len(docs), file_name)
        return docs
